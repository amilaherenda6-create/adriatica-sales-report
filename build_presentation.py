"""Build Adriatica Sales Management Presentation.
Run: python build_presentation.py
Output: adriatica_sales_report.pptx
"""
from pathlib import Path
import io
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Paths ──────────────────────────────────────────────────────────────────────
ROOT      = Path(__file__).parent
DATA      = ROOT / "data" / "orders.csv"
LOGO      = Path(__file__).parent.parent / "imb" / "images.png"
OUT       = ROOT / "adriatica_sales_report_v2.pptx"

# ── Brand colours ──────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1A, 0x35, 0x5E)   # IMB navy
MID_BLUE   = RGBColor(0x2E, 0x6D, 0xA4)   # accent
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)   # background tint
GOLD       = RGBColor(0xC8, 0x9F, 0x2F)   # highlight
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)
MID_GREY   = RGBColor(0x88, 0x88, 0x88)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Load & prepare data ────────────────────────────────────────────────────────
df = pd.read_csv(DATA, parse_dates=["order_date"])
df["revenue"]  = df["quantity"] * df["unit_price"]
df["month"]    = df["order_date"].dt.to_period("M").astype(str)
df["month_dt"] = pd.to_datetime(df["month"])
df["quarter"]  = df["order_date"].dt.to_period("Q").astype(str)

monthly    = df.groupby("month")["revenue"].sum().sort_index()
mom        = monthly.pct_change() * 100
cat_rev    = df.groupby("category")["revenue"].sum().sort_values(ascending=False)
top_prod   = df.groupby("product_name")["revenue"].sum().nlargest(10).sort_values()
top_cust   = df.groupby("customer_id").agg(
                revenue=("revenue","sum"), orders=("order_id","nunique")
             ).nlargest(10,"revenue").sort_values("revenue")

total_rev      = df["revenue"].sum()
best_month     = monthly.idxmax()
best_rev       = monthly.max()
avg_monthly    = monthly.mean()
unique_cust    = df["customer_id"].nunique()
total_orders   = df["order_id"].nunique()
last_mom       = mom.iloc[-1]
yoy_months     = list(monthly.index)
growth_h2h     = ((monthly.iloc[-1] - monthly.iloc[0]) / monthly.iloc[0]) * 100

# ── Chart helpers ──────────────────────────────────────────────────────────────
CHART_BG = "#F7FAFD"

def fig_to_stream(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=160, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def chart_monthly_revenue():
    fig, ax1 = plt.subplots(figsize=(11, 4), facecolor=CHART_BG)
    ax1.set_facecolor(CHART_BG)
    colors = ["#C8701A" if m == best_month else "#2E6DA4" for m in monthly.index]
    bars = ax1.bar(monthly.index, monthly.values / 1000, color=colors,
                   width=0.6, zorder=3)
    for bar, val in zip(bars, monthly.values):
        ax1.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.3,
                 f"€{val/1000:.0f}K", ha="center", va="bottom",
                 fontsize=7.5, color="#333333", fontweight="bold")
    ax2 = ax1.twinx()
    valid = mom.dropna()
    ax2.plot(valid.index, valid.values, color="#C8302F", linewidth=2,
             marker="o", markersize=5, zorder=4)
    ax2.axhline(0, color="#C8302F", linewidth=0.8, linestyle="--", alpha=0.4)
    ax1.set_ylabel("Revenue (€ thousands)", color="#1A355E", fontsize=9)
    ax2.set_ylabel("MoM Growth %", color="#C8302F", fontsize=9)
    ax1.tick_params(axis="x", rotation=35, labelsize=8)
    ax1.tick_params(axis="y", labelsize=8)
    ax2.tick_params(axis="y", labelsize=8, colors="#C8302F")
    ax1.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"€{x:.0f}K"))
    ax2.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{x:+.0f}%"))
    ax1.grid(axis="y", alpha=0.3, zorder=0)
    blue_patch  = mpatches.Patch(color="#2E6DA4", label="Monthly Revenue")
    gold_patch  = mpatches.Patch(color="#C8701A", label=f"Peak: {best_month}")
    red_line    = plt.Line2D([0],[0], color="#C8302F", linewidth=2, marker="o",
                              markersize=5, label="MoM Growth %")
    ax1.legend(handles=[blue_patch, gold_patch, red_line],
               loc="upper left", fontsize=8, framealpha=0.8)
    fig.tight_layout()
    return fig_to_stream(fig)


def chart_category():
    fig, ax = plt.subplots(figsize=(9, 4), facecolor=CHART_BG)
    ax.set_facecolor(CHART_BG)
    palette = ["#1A355E","#2E6DA4","#4A90C4","#6AAED6","#9ECAE1","#C6DBEF"]
    bars = ax.bar(cat_rev.index, cat_rev.values / 1000,
                  color=palette[:len(cat_rev)], width=0.6, zorder=3)
    for bar, val in zip(bars, cat_rev.values):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.4,
                f"€{val/1000:.0f}K", ha="center", va="bottom",
                fontsize=9, fontweight="bold", color="#1A355E")
    ax.set_ylabel("Revenue (€ thousands)", fontsize=9)
    ax.set_xlabel("")
    ax.tick_params(axis="x", labelsize=9)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"€{x:.0f}K"))
    ax.grid(axis="y", alpha=0.3, zorder=0)
    fig.tight_layout()
    return fig_to_stream(fig)


def chart_top_products():
    fig, ax = plt.subplots(figsize=(9, 5), facecolor=CHART_BG)
    ax.set_facecolor(CHART_BG)
    colors = ["#C8701A" if i == len(top_prod)-1 else "#2E6DA4"
              for i in range(len(top_prod))]
    bars = ax.barh(top_prod.index, top_prod.values / 1000,
                   color=colors, height=0.6, zorder=3)
    for bar, val in zip(bars, top_prod.values):
        ax.text(bar.get_width() + 0.2, bar.get_y() + bar.get_height()/2,
                f"€{val/1000:.1f}K", va="center", fontsize=8,
                fontweight="bold", color="#1A355E")
    ax.set_xlabel("Revenue (€ thousands)", fontsize=9)
    ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"€{x:.0f}K"))
    ax.tick_params(axis="y", labelsize=8)
    ax.grid(axis="x", alpha=0.3, zorder=0)
    fig.tight_layout()
    return fig_to_stream(fig)


def chart_top_customers():
    fig, ax = plt.subplots(figsize=(9, 5), facecolor=CHART_BG)
    ax.set_facecolor(CHART_BG)
    norm = plt.Normalize(top_cust["orders"].min(), top_cust["orders"].max())
    colors = plt.cm.Blues(norm(top_cust["orders"].values) * 0.7 + 0.3)
    bars = ax.barh(top_cust.index, top_cust["revenue"] / 1000,
                   color=colors, height=0.6, zorder=3)
    for bar, val, orders in zip(bars, top_cust["revenue"], top_cust["orders"]):
        ax.text(bar.get_width() + 0.1,
                bar.get_y() + bar.get_height()/2,
                f"€{val/1000:.1f}K  ({orders} orders)",
                va="center", fontsize=8, color="#1A355E")
    ax.set_xlabel("Revenue (€ thousands)", fontsize=9)
    ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"€{x:.0f}K"))
    ax.tick_params(axis="y", labelsize=8)
    ax.grid(axis="x", alpha=0.3, zorder=0)
    fig.tight_layout()
    return fig_to_stream(fig)


def chart_quarterly():
    q_rev = df.groupby("quarter")["revenue"].sum()
    fig, ax = plt.subplots(figsize=(7, 3.5), facecolor=CHART_BG)
    ax.set_facecolor(CHART_BG)
    colors = ["#C8701A" if q == "2024Q4" else "#2E6DA4" for q in q_rev.index]
    bars = ax.bar(q_rev.index, q_rev.values / 1000, color=colors,
                  width=0.5, zorder=3)
    for bar, val in zip(bars, q_rev.values):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.5,
                f"€{val/1000:.0f}K", ha="center", va="bottom",
                fontsize=9, fontweight="bold", color="#1A355E")
    ax.set_ylabel("Revenue (€ thousands)", fontsize=9)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"€{x:.0f}K"))
    ax.grid(axis="y", alpha=0.3, zorder=0)
    ax.tick_params(labelsize=9)
    fig.tight_layout()
    return fig_to_stream(fig)


# ── Slide builder helpers ──────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # completely blank


def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.width = line_width
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=Pt(12), bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def add_logo(slide, path, l, t, w):
    if path.exists():
        slide.shapes.add_picture(str(path), l, t, width=w)


def add_chart_image(slide, stream, l, t, w, h):
    slide.shapes.add_picture(stream, l, t, width=w, height=h)


def footer(slide, page_num, total):
    add_rect(slide, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35), DARK_BLUE)
    add_text(slide,
             "Adriatica Sales Report  |  Models in Advanced Data  |  Herenda Amila, IMB Master  |  22.06.2026",
             Inches(0.2), SLIDE_H - Inches(0.34), Inches(11.5), Inches(0.32),
             size=Pt(8.5), color=WHITE)
    add_text(slide, f"{page_num} / {total}",
             SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.34), Inches(1), Inches(0.32),
             size=Pt(8.5), color=WHITE, align=PP_ALIGN.RIGHT)


# ── Build slides ───────────────────────────────────────────────────────────────
prs   = new_prs()
TOTAL = 8

# ── 1. Title slide ─────────────────────────────────────────────────────────────
sl = blank_slide(prs)
# navy left panel
add_rect(sl, 0, 0, Inches(5.2), SLIDE_H, DARK_BLUE)
# white right panel (default)
# gold accent bar
add_rect(sl, Inches(5.2), 0, Inches(0.06), SLIDE_H, GOLD)

# logo on dark panel
add_logo(sl, LOGO, Inches(0.3), Inches(0.25), Inches(1.9))

# IMB label
add_text(sl, "International Master of Business",
         Inches(0.3), Inches(1.3), Inches(4.6), Inches(0.4),
         size=Pt(11), color=LIGHT_BLUE)
add_text(sl, "University of Ljubljana\nFaculty of Economics",
         Inches(0.3), Inches(1.68), Inches(4.6), Inches(0.65),
         size=Pt(10.5), color=WHITE)

# Gold divider line
add_rect(sl, Inches(0.3), Inches(2.45), Inches(4.5), Inches(0.04), GOLD)

add_text(sl, "ADRIATICA",
         Inches(0.3), Inches(2.6), Inches(4.6), Inches(0.85),
         size=Pt(38), bold=True, color=WHITE)
add_text(sl, "Sales Performance\nManagement Report",
         Inches(0.3), Inches(3.45), Inches(4.6), Inches(1.1),
         size=Pt(22), color=GOLD)
add_text(sl, "July 2024 – June 2025",
         Inches(0.3), Inches(4.62), Inches(4.6), Inches(0.38),
         size=Pt(13), color=LIGHT_BLUE)

# Course, author, date block
add_rect(sl, Inches(0.3), Inches(5.15), Inches(4.55), Inches(0.04), GOLD)
add_text(sl, "Course:  Models in Advanced Data",
         Inches(0.3), Inches(5.25), Inches(4.6), Inches(0.38),
         size=Pt(11), bold=True, color=WHITE)
add_text(sl, "Prepared by:  Herenda Amila, IMB Master",
         Inches(0.3), Inches(5.62), Inches(4.6), Inches(0.38),
         size=Pt(11), color=LIGHT_BLUE)
add_text(sl, "Date:  22.06.2026",
         Inches(0.3), Inches(5.98), Inches(4.6), Inches(0.38),
         size=Pt(11), color=LIGHT_BLUE)
add_text(sl, "CONFIDENTIAL  ·  FOR MANAGEMENT USE ONLY",
         Inches(0.3), Inches(6.55), Inches(4.6), Inches(0.3),
         size=Pt(8), color=MID_GREY)

# right panel — decorative KPI preview
kpis = [
    ("€822K",  "Total Revenue"),
    ("€97K",   "Peak Month (Dec 2024)"),
    ("1,093",  "Unique Customers"),
    ("4,644",  "Total Orders"),
]
for i, (val, lbl) in enumerate(kpis):
    top = Inches(1.5 + i * 1.3)
    add_rect(sl, Inches(5.6), top, Inches(7.4), Inches(1.1), LIGHT_BLUE)
    add_text(sl, val, Inches(5.8), top + Inches(0.08), Inches(4), Inches(0.55),
             size=Pt(28), bold=True, color=DARK_BLUE)
    add_text(sl, lbl, Inches(5.8), top + Inches(0.6), Inches(4), Inches(0.4),
             size=Pt(10), color=MID_BLUE)

footer(sl, 1, TOTAL)

# ── 2. Executive Summary ───────────────────────────────────────────────────────
sl = blank_slide(prs)
add_rect(sl, 0, 0, SLIDE_W, Inches(1.1), DARK_BLUE)
add_logo(sl, LOGO, Inches(0.15), Inches(0.1), Inches(0.8))
add_text(sl, "Executive Summary", Inches(1.1), Inches(0.2), Inches(10), Inches(0.7),
         size=Pt(30), bold=True, color=WHITE)
add_text(sl, "July 2024 – June 2025  ·  All categories",
         Inches(1.1), Inches(0.78), Inches(10), Inches(0.28),
         size=Pt(11), color=LIGHT_BLUE)

kpi_data = [
    ("€822,037", "Total Revenue", "12-month period"),
    ("€97,139",  "Peak Month",    "December 2024"),
    ("€68,503",  "Avg / Month",   "12-month average"),
    ("1,093",    "Unique Customers", "across all categories"),
]
col_w = Inches(3.0)
gap   = Inches(0.18)
for i, (val, title, sub) in enumerate(kpi_data):
    lx = Inches(0.25) + i * (col_w + gap)
    add_rect(sl, lx, Inches(1.25), col_w, Inches(1.8), LIGHT_BLUE)
    add_rect(sl, lx, Inches(1.25), col_w, Inches(0.08), MID_BLUE)
    add_text(sl, val,   lx+Inches(0.15), Inches(1.38), col_w-Inches(0.3), Inches(0.8),
             size=Pt(26), bold=True, color=DARK_BLUE)
    add_text(sl, title, lx+Inches(0.15), Inches(2.08), col_w-Inches(0.3), Inches(0.4),
             size=Pt(11), bold=True, color=DARK_BLUE)
    add_text(sl, sub,   lx+Inches(0.15), Inches(2.45), col_w-Inches(0.3), Inches(0.35),
             size=Pt(8.5), color=MID_GREY)

# Key findings
add_rect(sl, Inches(0.25), Inches(3.2), SLIDE_W - Inches(0.5), Inches(0.35), DARK_BLUE)
add_text(sl, "  KEY FINDINGS", Inches(0.3), Inches(3.22), Inches(6), Inches(0.3),
         size=Pt(10), bold=True, color=WHITE)

findings = [
    "▶  December 2024 was the strongest month at €97K — 42% above the annual average, driven by seasonal demand.",
    "▶  Electronics is the top-performing category, contributing the largest share of annual revenue.",
    "▶  The final month (June 2025) shows a MoM growth of {:.1f}%, indicating the trajectory at period end.".format(last_mom),
    "▶  Top 10 customers account for a disproportionate share of revenue — a concentration risk to monitor.",
    "▶  Total of {:,} orders placed by {:,} unique customers over the 12-month period.".format(total_orders, unique_cust),
]
for i, f in enumerate(findings):
    add_text(sl, f, Inches(0.35), Inches(3.65 + i * 0.6), SLIDE_W - Inches(0.7), Inches(0.55),
             size=Pt(10), color=DARK_GREY)

footer(sl, 2, TOTAL)

# ── 3. Monthly Revenue Trend ───────────────────────────────────────────────────
sl = blank_slide(prs)
add_rect(sl, 0, 0, SLIDE_W, Inches(1.1), DARK_BLUE)
add_logo(sl, LOGO, Inches(0.15), Inches(0.1), Inches(0.8))
add_text(sl, "Monthly Revenue & Growth Trend", Inches(1.1), Inches(0.2),
         Inches(11), Inches(0.7), size=Pt(30), bold=True, color=WHITE)
add_text(sl, "Orange bar = peak month (December 2024)  ·  Red line = Month-over-Month growth %",
         Inches(1.1), Inches(0.78), Inches(11), Inches(0.28), size=Pt(11), color=LIGHT_BLUE)

stream = chart_monthly_revenue()
add_chart_image(sl, stream, Inches(0.2), Inches(1.15), Inches(12.9), Inches(4.8))

# insight callout
add_rect(sl, Inches(0.25), Inches(6.1), Inches(12.8), Inches(1.0), LIGHT_BLUE)
add_text(sl,
         "Management Insight:  December 2024 revenue (€97K) was 42% above the 12-month average. "
         "Q4 holiday demand is the primary growth driver. Planning and inventory decisions "
         "should anticipate and capitalise on this seasonal pattern.",
         Inches(0.4), Inches(6.15), Inches(12.5), Inches(0.85),
         size=Pt(9.5), color=DARK_BLUE)

footer(sl, 3, TOTAL)

# ── 4. Quarterly View ──────────────────────────────────────────────────────────
sl = blank_slide(prs)
add_rect(sl, 0, 0, SLIDE_W, Inches(1.1), DARK_BLUE)
add_logo(sl, LOGO, Inches(0.15), Inches(0.1), Inches(0.8))
add_text(sl, "Quarterly Revenue Breakdown", Inches(1.1), Inches(0.2),
         Inches(11), Inches(0.7), size=Pt(30), bold=True, color=WHITE)
add_text(sl, "Q4 2024 highlighted  ·  Consistent base with strong seasonal uplift",
         Inches(1.1), Inches(0.78), Inches(11), Inches(0.28), size=Pt(11), color=LIGHT_BLUE)

stream_q = chart_quarterly()
add_chart_image(sl, stream_q, Inches(0.3), Inches(1.2), Inches(7.2), Inches(4.5))

# Q commentary
q_rev = df.groupby("quarter")["revenue"].sum()
bullets = []
for q, v in q_rev.items():
    share = v / total_rev * 100
    bullets.append(f"  {q}   €{v:,.0f}   ({share:.1f}% of annual revenue)")

add_rect(sl, Inches(7.8), Inches(1.2), Inches(5.2), Inches(4.5), LIGHT_BLUE)
add_text(sl, "Quarterly Breakdown", Inches(8.0), Inches(1.3), Inches(4.8), Inches(0.4),
         size=Pt(11), bold=True, color=DARK_BLUE)
for i, b in enumerate(bullets):
    add_text(sl, b, Inches(8.0), Inches(1.85 + i * 0.65), Inches(4.8), Inches(0.55),
             size=Pt(10.5), color=DARK_GREY)

add_text(sl,
         "Q4 (Oct–Dec 2024) is the standout quarter. "
         "Management should allocate sales resource and marketing spend "
         "to maximise Q4 capture in the next planning cycle.",
         Inches(7.8), Inches(5.9), Inches(5.2), Inches(0.85),
         size=Pt(9), color=MID_BLUE)

footer(sl, 4, TOTAL)

# ── 5. Category Performance ────────────────────────────────────────────────────
sl = blank_slide(prs)
add_rect(sl, 0, 0, SLIDE_W, Inches(1.1), DARK_BLUE)
add_logo(sl, LOGO, Inches(0.15), Inches(0.1), Inches(0.8))
add_text(sl, "Revenue by Category", Inches(1.1), Inches(0.2),
         Inches(11), Inches(0.7), size=Pt(30), bold=True, color=WHITE)
add_text(sl, "Which categories drive the business — and which have headroom to grow",
         Inches(1.1), Inches(0.78), Inches(11), Inches(0.28), size=Pt(11), color=LIGHT_BLUE)

stream_cat = chart_category()
add_chart_image(sl, stream_cat, Inches(0.2), Inches(1.15), Inches(9.2), Inches(4.5))

# Category table
add_rect(sl, Inches(9.5), Inches(1.2), Inches(3.6), Inches(4.5), LIGHT_BLUE)
add_text(sl, "Category  |  Revenue  |  Share",
         Inches(9.6), Inches(1.28), Inches(3.4), Inches(0.3),
         size=Pt(8), bold=True, color=DARK_BLUE)
for i, (cat, rev) in enumerate(cat_rev.items()):
    share = rev / total_rev * 100
    add_rect(sl, Inches(9.5), Inches(1.6 + i * 0.62), Inches(3.6), Inches(0.6),
             WHITE if i % 2 == 0 else None)
    add_text(sl, cat, Inches(9.6), Inches(1.65 + i * 0.62), Inches(1.5), Inches(0.5),
             size=Pt(9), color=DARK_GREY)
    add_text(sl, f"€{rev:,.0f}", Inches(11.1), Inches(1.65 + i * 0.62),
             Inches(1.1), Inches(0.5), size=Pt(9), bold=True, color=DARK_BLUE)
    add_text(sl, f"{share:.1f}%", Inches(12.2), Inches(1.65 + i * 0.62),
             Inches(0.8), Inches(0.5), size=Pt(9), color=MID_BLUE)

add_rect(sl, Inches(0.25), Inches(5.85), Inches(12.8), Inches(0.85), LIGHT_BLUE)
add_text(sl,
         "Management Insight:  The top category drives disproportionate revenue. "
         "Diversification into mid-tier categories (Office, Home & Kitchen) "
         "reduces concentration risk and opens new growth vectors.",
         Inches(0.4), Inches(5.9), Inches(12.5), Inches(0.75),
         size=Pt(9.5), color=DARK_BLUE)

footer(sl, 5, TOTAL)

# ── 6. Top 10 Products ────────────────────────────────────────────────────────
sl = blank_slide(prs)
add_rect(sl, 0, 0, SLIDE_W, Inches(1.1), DARK_BLUE)
add_logo(sl, LOGO, Inches(0.15), Inches(0.1), Inches(0.8))
add_text(sl, "Top 10 Products by Revenue", Inches(1.1), Inches(0.2),
         Inches(11), Inches(0.7), size=Pt(30), bold=True, color=WHITE)
add_text(sl, "Orange bar = #1 product  ·  Focus assortment planning on these revenue anchors",
         Inches(1.1), Inches(0.78), Inches(11), Inches(0.28), size=Pt(11), color=LIGHT_BLUE)

stream_prod = chart_top_products()
add_chart_image(sl, stream_prod, Inches(0.2), Inches(1.15), Inches(12.9), Inches(5.0))

add_rect(sl, Inches(0.25), Inches(6.25), Inches(12.8), Inches(0.85), LIGHT_BLUE)
top_p_name = top_prod.index[-1]
top_p_rev  = top_prod.values[-1]
add_text(sl,
         f"Management Insight:  '{top_p_name}' leads with €{top_p_rev:,.0f}. "
         "The top 10 products represent a concentrated revenue engine — "
         "stock availability and promotional priority for these SKUs directly impacts monthly performance.",
         Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.75),
         size=Pt(9.5), color=DARK_BLUE)

footer(sl, 6, TOTAL)

# ── 7. Top 10 Customers ───────────────────────────────────────────────────────
sl = blank_slide(prs)
add_rect(sl, 0, 0, SLIDE_W, Inches(1.1), DARK_BLUE)
add_logo(sl, LOGO, Inches(0.15), Inches(0.1), Inches(0.8))
add_text(sl, "Top 10 Customers by Revenue", Inches(1.1), Inches(0.2),
         Inches(11), Inches(0.7), size=Pt(30), bold=True, color=WHITE)
add_text(sl, "Bar colour = order frequency  ·  Darker = more orders  ·  Brackets show order count",
         Inches(1.1), Inches(0.78), Inches(11), Inches(0.28), size=Pt(11), color=LIGHT_BLUE)

stream_cust = chart_top_customers()
add_chart_image(sl, stream_cust, Inches(0.2), Inches(1.15), Inches(12.9), Inches(5.0))

add_rect(sl, Inches(0.25), Inches(6.25), Inches(12.8), Inches(0.85), LIGHT_BLUE)
top_c_id  = top_cust.index[-1]
top_c_rev = top_cust["revenue"].values[-1]
add_text(sl,
         f"Management Insight:  Customer {top_c_id} is the highest-value account (€{top_c_rev:,.0f}). "
         "High-value, high-frequency customers warrant dedicated account management. "
         "Retention of the top 10 is a strategic priority — their churn would materially impact annual revenue.",
         Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.75),
         size=Pt(9.5), color=DARK_BLUE)

footer(sl, 7, TOTAL)

# ── 8. Key Takeaways & Next Steps ─────────────────────────────────────────────
sl = blank_slide(prs)
add_rect(sl, 0, 0, SLIDE_W, Inches(1.1), DARK_BLUE)
add_logo(sl, LOGO, Inches(0.15), Inches(0.1), Inches(0.8))
add_text(sl, "Key Takeaways & Management Actions", Inches(1.1), Inches(0.2),
         Inches(11), Inches(0.7), size=Pt(30), bold=True, color=WHITE)
add_text(sl, "Strategic recommendations based on 12-month performance data",
         Inches(1.1), Inches(0.78), Inches(11), Inches(0.28), size=Pt(11), color=LIGHT_BLUE)

takeaways = [
    ("1", "Capitalise on Q4 seasonality",
     "December revenue was 42% above the monthly average. "
     "Increase inventory, marketing spend, and sales capacity in Q4 to maximise this structural advantage."),
    ("2", "Protect & grow top accounts",
     "Top 10 customers generate an outsized share of revenue. "
     "Assign dedicated account managers, offer loyalty incentives, and monitor churn risk monthly."),
    ("3", "Grow mid-tier categories",
     "Office, Home & Kitchen, and Toys are under-indexed relative to Electronics. "
     "A targeted category development plan could unlock €50K–100K in incremental revenue."),
    ("4", "Strengthen top-10 product availability",
     "SKU concentration means stock-outs directly hit revenue. "
     "Implement safety-stock rules and lead-time monitoring for the top 10 products."),
    ("5", "Set MoM growth targets",
     "Month-over-Month growth is volatile (-20% to +40%). "
     "Establish monthly revenue targets and an early-warning dashboard to catch shortfalls before quarter end."),
]

for i, (num, title, body) in enumerate(takeaways):
    ty = Inches(1.25 + i * 1.17)
    add_rect(sl, Inches(0.25), ty, Inches(0.55), Inches(1.0), GOLD)
    add_text(sl, num, Inches(0.25), ty + Inches(0.15), Inches(0.55), Inches(0.7),
             size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, Inches(0.85), ty, Inches(12.2), Inches(1.0), LIGHT_BLUE)
    add_text(sl, title, Inches(1.0), ty + Inches(0.05), Inches(12.0), Inches(0.4),
             size=Pt(11), bold=True, color=DARK_BLUE)
    add_text(sl, body,  Inches(1.0), ty + Inches(0.45), Inches(12.0), Inches(0.5),
             size=Pt(9.5), color=DARK_GREY)

footer(sl, 8, TOTAL)

# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"Saved: {OUT}")
